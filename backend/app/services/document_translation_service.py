"""
Document Translation Service - Universal Pivot Workflow

This service implements format-preserving translation by:
1. Converting binary formats (PDF) to editable formats (DOCX)
2. Extracting text content into a structured JSON format
3. Translating the JSON content via LLM
4. Re-injecting translated text back into the document
5. Converting back to original format if needed

Supports: PDF, DOCX, PPTX with layout preservation.
"""

import atexit
import json
import logging
import os
import re
import shutil
import tempfile
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Any, Optional
from copy import deepcopy

import structlog

logger = structlog.get_logger()


@dataclass
class TextSegment:
    """A segment of text with its styling and position."""
    id: str
    text: str
    # Styling
    font_name: Optional[str] = None
    font_size: Optional[float] = None
    bold: bool = False
    italic: bool = False
    underline: bool = False
    color: Optional[str] = None
    # Position context
    paragraph_index: int = 0
    run_index: int = 0
    # For tables
    table_index: Optional[int] = None
    row_index: Optional[int] = None
    cell_index: Optional[int] = None
    # For slides
    slide_index: Optional[int] = None
    shape_index: Optional[int] = None


@dataclass
class DocumentStructure:
    """Structured representation of a document for translation."""
    source_file: str
    source_format: str
    segments: list[TextSegment] = field(default_factory=list)
    metadata: dict = field(default_factory=dict)

    def to_json(self) -> str:
        """Serialize to JSON for translation."""
        return json.dumps({
            "source_file": self.source_file,
            "source_format": self.source_format,
            "segments": [asdict(s) for s in self.segments],
            "metadata": self.metadata,
        }, indent=2, ensure_ascii=False)

    @classmethod
    def from_json(cls, json_str: str) -> "DocumentStructure":
        """Deserialize from JSON."""
        data = json.loads(json_str)
        segments = [TextSegment(**s) for s in data.get("segments", [])]
        return cls(
            source_file=data["source_file"],
            source_format=data["source_format"],
            segments=segments,
            metadata=data.get("metadata", {}),
        )

    def get_translation_payload(self) -> list[dict]:
        """Get segments formatted for translation API."""
        return [
            {"id": s.id, "text": s.text}
            for s in self.segments
            if s.text.strip()  # Skip empty segments
        ]

    def apply_translations(self, translations: dict[str, str]):
        """Apply translated text back to segments."""
        for segment in self.segments:
            if segment.id in translations:
                segment.text = translations[segment.id]


class DocumentTranslationService:
    """
    Service for format-preserving document translation.

    Workflow:
    1. PDF → DOCX (if needed) using pdf2docx
    2. Extract text segments to JSON structure
    3. Translate segments via LLM
    4. Re-inject translated text with font scaling
    5. Convert back to original format
    """

    # Configuration constants for font scaling
    FONT_SCALE_THRESHOLD = 1.25  # Apply scaling when text expands by 25%
    MIN_FONT_SIZE_PT = 8  # Minimum font size in points
    MAX_SCALE_FACTOR = 0.85  # Maximum reduction (85% of original)
    
    # Configuration for translation validation
    TRANSLATION_COMPLETENESS_THRESHOLD = 0.9  # Warn if less than 90% of translations applied
    
    def __init__(self):
        self._temp_dir = tempfile.mkdtemp(prefix="doc_translate_")
        self._job_temp_dirs = []  # Track per-job temp directories

    # =========================================================================
    # PDF to DOCX Conversion
    # =========================================================================

    async def pdf_to_docx(self, pdf_path: str, output_path: Optional[str] = None) -> str:
        """
        Convert PDF to DOCX using pdf2docx.

        This preserves:
        - Text with formatting
        - Tables
        - Images
        - Basic layout
        """
        from pdf2docx import Converter

        if output_path is None:
            output_path = os.path.join(
                self._temp_dir,
                Path(pdf_path).stem + ".docx"
            )

        logger.info("Converting PDF to DOCX", pdf_path=pdf_path, output_path=output_path)

        cv = Converter(pdf_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()

        logger.info("PDF to DOCX conversion complete", output_path=output_path)
        return output_path

    async def docx_to_pdf(self, docx_path: str, output_path: Optional[str] = None) -> str:
        """
        Convert DOCX back to PDF.

        Uses LibreOffice in headless mode or alternative.
        """
        if output_path is None:
            output_path = os.path.join(
                self._temp_dir,
                Path(docx_path).stem + ".pdf"
            )

        # Try using LibreOffice
        try:
            import subprocess
            result = subprocess.run([
                "libreoffice", "--headless", "--convert-to", "pdf",
                "--outdir", os.path.dirname(output_path),
                docx_path
            ], capture_output=True, timeout=120)

            if result.returncode == 0:
                # LibreOffice outputs to same dir with .pdf extension
                expected_output = os.path.join(
                    os.path.dirname(output_path),
                    Path(docx_path).stem + ".pdf"
                )
                if os.path.exists(expected_output) and expected_output != output_path:
                    os.rename(expected_output, output_path)
                return output_path
        except Exception as e:
            logger.warning("LibreOffice conversion failed", error=str(e))

        # Fallback: return DOCX (user can convert manually)
        logger.warning("PDF conversion not available, returning DOCX")
        return docx_path

    # =========================================================================
    # DOCX Text Extraction & Re-injection
    # =========================================================================

    async def extract_docx_structure(self, docx_path: str) -> DocumentStructure:
        """
        Extract text structure from DOCX file.

        Preserves paragraph and run structure for accurate re-injection.
        """
        from docx import Document
        from docx.shared import Pt, RGBColor

        doc = Document(docx_path)
        structure = DocumentStructure(
            source_file=docx_path,
            source_format="docx",
            metadata={"paragraph_count": len(doc.paragraphs)}
        )

        segment_id = 0

        # Extract paragraphs
        for para_idx, paragraph in enumerate(doc.paragraphs):
            for run_idx, run in enumerate(paragraph.runs):
                if run.text.strip():
                    segment = TextSegment(
                        id=f"p{para_idx}_r{run_idx}",
                        text=run.text,
                        paragraph_index=para_idx,
                        run_index=run_idx,
                        font_name=run.font.name,
                        font_size=run.font.size.pt if run.font.size else None,
                        bold=run.bold or False,
                        italic=run.italic or False,
                        underline=run.underline or False,
                    )
                    structure.segments.append(segment)
                    segment_id += 1

        # Extract tables
        for table_idx, table in enumerate(doc.tables):
            for row_idx, row in enumerate(table.rows):
                for cell_idx, cell in enumerate(row.cells):
                    for para_idx, paragraph in enumerate(cell.paragraphs):
                        for run_idx, run in enumerate(paragraph.runs):
                            if run.text.strip():
                                segment = TextSegment(
                                    id=f"t{table_idx}_r{row_idx}_c{cell_idx}_p{para_idx}_r{run_idx}",
                                    text=run.text,
                                    table_index=table_idx,
                                    row_index=row_idx,
                                    cell_index=cell_idx,
                                    paragraph_index=para_idx,
                                    run_index=run_idx,
                                    font_name=run.font.name,
                                    font_size=run.font.size.pt if run.font.size else None,
                                    bold=run.bold or False,
                                    italic=run.italic or False,
                                )
                                structure.segments.append(segment)
                                segment_id += 1

        logger.info(
            "Extracted DOCX structure",
            segments=len(structure.segments),
            paragraphs=len(doc.paragraphs),
            tables=len(doc.tables),
        )

        return structure

    async def inject_docx_translations(
        self,
        docx_path: str,
        structure: DocumentStructure,
        output_path: str,
        scale_fonts: bool = True,
    ) -> str:
        """
        Re-inject translated text into DOCX while preserving formatting.

        Implements dynamic font scaling for content that expands.
        """
        from docx import Document
        from docx.shared import Pt

        # Copy original file
        shutil.copy(docx_path, output_path)
        doc = Document(output_path)

        # Build lookup for translations
        translations = {s.id: s for s in structure.segments}
        
        # Track translation statistics for validation
        applied_count = 0
        skipped_count = 0

        # Inject into paragraphs
        for para_idx, paragraph in enumerate(doc.paragraphs):
            for run_idx, run in enumerate(paragraph.runs):
                seg_id = f"p{para_idx}_r{run_idx}"
                if seg_id in translations:
                    segment = translations[seg_id]
                    original_len = len(run.text)
                    translated_len = len(segment.text)

                    run.text = segment.text
                    
                    # Re-apply formatting attributes that were preserved
                    if segment.bold is not None:
                        run.bold = segment.bold
                    if segment.italic is not None:
                        run.italic = segment.italic
                    if segment.underline is not None:
                        run.underline = segment.underline
                    
                    applied_count += 1

                    # Dynamic font scaling if text expanded significantly
                    # Use consistent threshold across all content types
                    if scale_fonts and translated_len > original_len * self.FONT_SCALE_THRESHOLD:
                        if run.font.size:
                            scale_factor = min(original_len / translated_len, self.MAX_SCALE_FACTOR)
                            new_size = max(run.font.size.pt * scale_factor, self.MIN_FONT_SIZE_PT)
                            run.font.size = Pt(new_size)
                else:
                    # Track segments that weren't translated
                    if run.text.strip():
                        skipped_count += 1

        # Inject into tables
        for table_idx, table in enumerate(doc.tables):
            for row_idx, row in enumerate(table.rows):
                for cell_idx, cell in enumerate(row.cells):
                    for para_idx, paragraph in enumerate(cell.paragraphs):
                        for run_idx, run in enumerate(paragraph.runs):
                            seg_id = f"t{table_idx}_r{row_idx}_c{cell_idx}_p{para_idx}_r{run_idx}"
                            if seg_id in translations:
                                segment = translations[seg_id]
                                original_len = len(run.text)
                                translated_len = len(segment.text)

                                run.text = segment.text
                                
                                # Re-apply formatting attributes
                                if segment.bold is not None:
                                    run.bold = segment.bold
                                if segment.italic is not None:
                                    run.italic = segment.italic
                                if segment.underline is not None:
                                    run.underline = segment.underline
                                
                                applied_count += 1

                                # Use consistent scaling threshold for tables
                                if scale_fonts and translated_len > original_len * self.FONT_SCALE_THRESHOLD:
                                    if run.font.size:
                                        scale_factor = min(original_len / translated_len, self.MAX_SCALE_FACTOR)
                                        new_size = max(run.font.size.pt * scale_factor, self.MIN_FONT_SIZE_PT)
                                        run.font.size = Pt(new_size)
                            else:
                                if run.text.strip():
                                    skipped_count += 1

        doc.save(output_path)
        
        # Log translation injection statistics
        total_expected = len(translations)
        logger.info(
            "Injected translations into DOCX",
            output_path=output_path,
            applied=applied_count,
            expected=total_expected,
            skipped=skipped_count,
        )
        
        # Warn if significant mismatch
        if total_expected > 0 and applied_count < total_expected * self.TRANSLATION_COMPLETENESS_THRESHOLD:
            logger.warning(
                "Translation injection incomplete",
                applied=applied_count,
                expected=total_expected,
                missing_pct=int((1 - applied_count/total_expected) * 100)
            )

        return output_path

    # =========================================================================
    # PPTX Text Extraction & Re-injection
    # =========================================================================

    async def extract_pptx_structure(self, pptx_path: str) -> DocumentStructure:
        """Extract text structure from PowerPoint file."""
        from pptx import Presentation

        prs = Presentation(pptx_path)
        structure = DocumentStructure(
            source_file=pptx_path,
            source_format="pptx",
            metadata={"slide_count": len(prs.slides)}
        )

        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if not shape.has_text_frame:
                    continue

                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    for run_idx, run in enumerate(paragraph.runs):
                        if run.text.strip():
                            segment = TextSegment(
                                id=f"s{slide_idx}_sh{shape_idx}_p{para_idx}_r{run_idx}",
                                text=run.text,
                                slide_index=slide_idx,
                                shape_index=shape_idx,
                                paragraph_index=para_idx,
                                run_index=run_idx,
                                font_name=run.font.name,
                                font_size=run.font.size.pt if run.font.size else None,
                                bold=run.font.bold or False,
                                italic=run.font.italic or False,
                            )
                            structure.segments.append(segment)

        logger.info(
            "Extracted PPTX structure",
            segments=len(structure.segments),
            slides=len(prs.slides),
        )

        return structure

    async def inject_pptx_translations(
        self,
        pptx_path: str,
        structure: DocumentStructure,
        output_path: str,
        scale_fonts: bool = True,
    ) -> str:
        """Re-inject translated text into PPTX with font scaling."""
        from pptx import Presentation
        from pptx.util import Pt

        shutil.copy(pptx_path, output_path)
        prs = Presentation(output_path)

        translations = {s.id: s for s in structure.segments}
        
        # Track translation statistics
        applied_count = 0
        skipped_count = 0

        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if not shape.has_text_frame:
                    continue

                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    for run_idx, run in enumerate(paragraph.runs):
                        seg_id = f"s{slide_idx}_sh{shape_idx}_p{para_idx}_r{run_idx}"
                        if seg_id in translations:
                            segment = translations[seg_id]
                            original_len = len(run.text)
                            translated_len = len(segment.text)

                            run.text = segment.text
                            
                            # Re-apply formatting attributes
                            if segment.bold is not None:
                                run.font.bold = segment.bold
                            if segment.italic is not None:
                                run.font.italic = segment.italic
                            
                            applied_count += 1

                            # Use consistent scaling threshold for presentations
                            if scale_fonts and translated_len > original_len * self.FONT_SCALE_THRESHOLD:
                                if run.font.size:
                                    scale_factor = min(original_len / translated_len, self.MAX_SCALE_FACTOR)
                                    new_size = max(run.font.size.pt * scale_factor, self.MIN_FONT_SIZE_PT)
                                    run.font.size = Pt(new_size)
                        else:
                            if run.text.strip():
                                skipped_count += 1

        prs.save(output_path)
        
        # Log translation injection statistics
        total_expected = len(translations)
        logger.info(
            "Injected translations into PPTX",
            output_path=output_path,
            applied=applied_count,
            expected=total_expected,
            skipped=skipped_count,
        )
        
        # Warn if significant mismatch
        if total_expected > 0 and applied_count < total_expected * self.TRANSLATION_COMPLETENESS_THRESHOLD:
            logger.warning(
                "Translation injection incomplete",
                applied=applied_count,
                expected=total_expected,
                missing_pct=int((1 - applied_count/total_expected) * 100)
            )

        return output_path

    # =========================================================================
    # Translation via LLM
    # =========================================================================

    async def translate_segments(
        self,
        segments: list[dict],
        target_language: str,
        simplified: bool,
        llm_model: str,
        agent=None,
        progress_callback=None,
    ) -> dict[str, str]:
        """
        Translate text segments using LLM.

        Args:
            segments: List of segments with 'id' and 'text' keys
            target_language: Target language code
            simplified: Use simplified language
            llm_model: LLM model name
            agent: Agent object
            progress_callback: Optional async callback(completed, total) for progress updates

        Returns a dict mapping segment IDs to translated text.
        """
        from .translation_service import translation_service

        # Build translation prompt
        translations = {}

        # Process in batches to avoid token limits
        # Increased from 50 to 200 to leverage the 128K context window of gpt-oss-20b
        # This reduces API calls by 4x and speeds up translation significantly
        batch_size = 200
        total_batches = (len(segments) + batch_size - 1) // batch_size

        for batch_num, i in enumerate(range(0, len(segments), batch_size)):
            batch = segments[i:i + batch_size]

            # Format segments for translation
            segments_text = "\n".join([
                f"[{s['id']}]: {s['text']}"
                for s in batch
            ])

            # Create translation prompt
            prompt = self._build_translation_prompt(
                segments_text, target_language, simplified
            )

            # Call LLM
            result = await translation_service.translate_text(
                text=prompt,
                target_language=target_language,
                agent=agent,
                simplified=simplified,
            )

            # Parse response
            translated_text = result.get("translated_text", "")
            batch_translations = self._parse_translation_response(translated_text)
            
            # Validate translation completeness for this batch
            expected_ids = {s['id'] for s in batch}
            found_ids = set(batch_translations.keys())
            missing_ids = expected_ids - found_ids
            
            if missing_ids:
                logger.warning(
                    f"Batch {batch_num + 1}/{total_batches} missing translations",
                    missing_count=len(missing_ids),
                    missing_ids=list(missing_ids)[:5],  # Log first 5
                )
                
                # For missing segments, use original text as fallback
                for seg in batch:
                    if seg['id'] in missing_ids:
                        logger.debug(f"Using original text as fallback for {seg['id']}")
                        batch_translations[seg['id']] = seg['text']
            
            translations.update(batch_translations)

            # Report progress
            if progress_callback:
                try:
                    await progress_callback(batch_num + 1, total_batches)
                except Exception as e:
                    logger.warning(f"Progress callback failed: {e}")

            logger.info(
                f"Translated batch {batch_num + 1}/{total_batches}",
                segments_in_batch=len(batch),
                translations_found=len(batch_translations),
                missing_count=len(missing_ids),
            )
        
        # Final validation
        total_segments = len(segments)
        total_translations = len(translations)
        if total_translations < total_segments:
            logger.warning(
                "Translation incomplete",
                expected=total_segments,
                received=total_translations,
                missing=total_segments - total_translations,
            )

        return translations

    def _build_translation_prompt(
        self,
        segments_text: str,
        target_language: str,
        simplified: bool,
    ) -> str:
        """Build the prompt for segment translation."""
        mode = "simplified, easy-to-understand language" if simplified else "standard language"

        return f"""Translate the following text segments to {target_language} using {mode}.

IMPORTANT RULES:
1. Maintain the EXACT format: [segment_id]: translated_text
2. Each segment must be on its own line
3. Do not add or remove segments - translate ALL segments provided
4. Preserve any formatting markers, numbers, and special characters
5. Keep proper nouns, brand names, and technical terms as appropriate
6. Do NOT add explanations, notes, or any additional commentary
7. Output ONLY the translated segments in the required format

TEXT SEGMENTS TO TRANSLATE:
{segments_text}

TRANSLATED SEGMENTS:"""

    def _parse_translation_response(self, response: str) -> dict[str, str]:
        """
        Parse the LLM response to extract segment translations.
        
        Uses multiple parsing strategies for robustness.
        """
        translations = {}

        # Strategy 1: Match pattern [segment_id]: translated_text
        pattern = r'\[([^\]]+)\]:\s*(.+?)(?=\n\[|\Z)'
        matches = re.findall(pattern, response, re.DOTALL)
        
        for seg_id, text in matches:
            translations[seg_id.strip()] = text.strip()
        
        # Strategy 2: If primary parsing found nothing, try line-by-line fallback
        # This handles cases where the LLM didn't follow the exact format
        if not translations and response:
            logger.warning("Primary segment parsing failed, attempting fallback parsing")
            for line in response.split('\n'):
                # Try to match [id]: text pattern more leniently
                if ']:' in line:
                    try:
                        # Find the first ] followed by :
                        bracket_close = line.index(']')
                        # Check if next character is a colon (with bounds checking)
                        if bracket_close + 1 < len(line) and line[bracket_close + 1] == ':':
                            seg_id = line[1:bracket_close].strip()  # Skip opening [
                            text = line[bracket_close+2:].strip()
                            if seg_id and text:
                                translations[seg_id] = text
                    except (ValueError, IndexError):
                        continue
        
        return translations

    # =========================================================================
    # Main Translation Pipeline
    # =========================================================================

    async def translate_document(
        self,
        input_path: str,
        output_path: str,
        target_language: str,
        simplified: bool = False,
        llm_model: str = "gpt-4",
        agent=None,
        preserve_format: bool = True,
    ) -> dict[str, Any]:
        """
        Main entry point for format-preserving document translation.

        Args:
            input_path: Path to source document
            output_path: Path for translated document
            target_language: Target language code
            simplified: Use simplified language
            llm_model: LLM model for translation
            agent: Translation agent (optional)
            preserve_format: Output in same format as input

        Returns:
            dict with translation stats and output path
        """
        input_ext = Path(input_path).suffix.lower()

        logger.info(
            "Starting document translation",
            input_path=input_path,
            format=input_ext,
            target_language=target_language,
        )

        # Determine workflow based on input format
        if input_ext == ".pdf":
            return await self._translate_pdf(
                input_path, output_path, target_language,
                simplified, llm_model, agent, preserve_format
            )
        elif input_ext in (".docx", ".doc"):
            return await self._translate_docx(
                input_path, output_path, target_language,
                simplified, llm_model, agent
            )
        elif input_ext in (".pptx", ".ppt"):
            return await self._translate_pptx(
                input_path, output_path, target_language,
                simplified, llm_model, agent
            )
        else:
            raise ValueError(f"Unsupported format: {input_ext}")

    async def _translate_pdf(
        self,
        input_path: str,
        output_path: str,
        target_language: str,
        simplified: bool,
        llm_model: str,
        agent,
        preserve_format: bool,
    ) -> dict[str, Any]:
        """Translate PDF via DOCX pivot."""

        # Step 1: Convert PDF to DOCX
        docx_path = await self.pdf_to_docx(input_path)

        # Step 2: Translate DOCX
        translated_docx = os.path.join(
            self._temp_dir,
            Path(input_path).stem + f"_{target_language}.docx"
        )
        result = await self._translate_docx(
            docx_path, translated_docx, target_language,
            simplified, llm_model, agent
        )

        # Step 3: Convert back to PDF if requested
        if preserve_format and output_path.endswith(".pdf"):
            final_path = await self.docx_to_pdf(translated_docx, output_path)
        else:
            # Output as DOCX
            final_output = output_path.replace(".pdf", ".docx")
            shutil.copy(translated_docx, final_output)
            final_path = final_output

        result["output_path"] = final_path
        result["output_format"] = Path(final_path).suffix

        return result

    async def _translate_docx(
        self,
        input_path: str,
        output_path: str,
        target_language: str,
        simplified: bool,
        llm_model: str,
        agent,
    ) -> dict[str, Any]:
        """Translate DOCX file."""

        # Extract structure
        structure = await self.extract_docx_structure(input_path)

        # Get segments for translation
        segments = structure.get_translation_payload()

        if not segments:
            raise ValueError("No translatable content found in document")

        # Translate segments
        translations = await self.translate_segments(
            segments, target_language, simplified, llm_model, agent
        )

        # Apply translations
        structure.apply_translations(translations)

        # Re-inject into document
        await self.inject_docx_translations(
            input_path, structure, output_path, scale_fonts=True
        )

        return {
            "output_path": output_path,
            "output_format": ".docx",
            "segments_translated": len(translations),
            "total_segments": len(segments),
        }

    async def _translate_pptx(
        self,
        input_path: str,
        output_path: str,
        target_language: str,
        simplified: bool,
        llm_model: str,
        agent,
    ) -> dict[str, Any]:
        """Translate PowerPoint file."""

        # Extract structure
        structure = await self.extract_pptx_structure(input_path)

        # Get segments for translation
        segments = structure.get_translation_payload()

        if not segments:
            raise ValueError("No translatable content found in presentation")

        # Translate segments
        translations = await self.translate_segments(
            segments, target_language, simplified, llm_model, agent
        )

        # Apply translations
        structure.apply_translations(translations)

        # Re-inject into document
        await self.inject_pptx_translations(
            input_path, structure, output_path, scale_fonts=True
        )

        return {
            "output_path": output_path,
            "output_format": ".pptx",
            "segments_translated": len(translations),
            "total_segments": len(segments),
        }

    def create_job_temp_dir(self) -> str:
        """
        Create a job-specific temporary directory for better isolation.
        
        Returns the path to the new temp directory.
        """
        job_temp = tempfile.mkdtemp(prefix="doc_job_", dir=self._temp_dir)
        self._job_temp_dirs.append(job_temp)
        return job_temp
    
    def cleanup_job_temp_dir(self, job_temp_dir: str):
        """Clean up a specific job's temporary directory."""
        if os.path.exists(job_temp_dir):
            try:
                shutil.rmtree(job_temp_dir)
                if job_temp_dir in self._job_temp_dirs:
                    self._job_temp_dirs.remove(job_temp_dir)
                logger.debug("Cleaned up job temp directory", path=job_temp_dir)
            except Exception as e:
                logger.warning("Failed to clean up job temp directory", path=job_temp_dir, error=str(e))
    
    def cleanup(self):
        """Clean up all temporary files and directories."""
        
        # Clean up individual job directories first
        for job_dir in list(self._job_temp_dirs):
            self.cleanup_job_temp_dir(job_dir)
        
        # Clean up main temp directory
        if os.path.exists(self._temp_dir):
            try:
                shutil.rmtree(self._temp_dir)
                logger.debug("Cleaned up main temp directory", path=self._temp_dir)
            except Exception as e:
                logger.warning("Failed to clean up main temp directory", path=self._temp_dir, error=str(e))
    
    def __enter__(self):
        """Context manager entry."""
        return self
    
    def __exit__(self, exc_type, exc_val, exc_tb):
        """Context manager exit - ensures cleanup."""
        self.cleanup()
        return False


# Singleton instance
document_translation_service = DocumentTranslationService()

# Register cleanup on exit to ensure temp files are removed
atexit.register(document_translation_service.cleanup)

